from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import List
import os
from datetime import datetime
import re

def parse_grok_grouped_lyrics(lyrics: str) -> List[List[str]]:
    """
    Parse lyrics that Grok has already grouped with ---SLIDE--- separators
    """
    # Split by the slide separator
    sections = lyrics.split('---SLIDE---')

    groups = []
    for section in sections:
        # Split into lines and filter out empty ones
        lines = [line.strip() for line in section.strip().split('\n') if line.strip()]
        if lines:
            groups.append(lines)

    return groups

def create_title_slide(prs: Presentation, song_name: str):
    """
    Create a title slide for the song
    """
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add title text box (centered)
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(2)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = 1  # Middle alignment

    p = text_frame.paragraphs[0]
    p.text = song_name
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.name = 'Arial'
    p.font.color.rgb = RGBColor(0, 0, 0)

    # Set slide background to white
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

def create_lyrics_slide(prs: Presentation, lyrics_chunk: List[str]):
    """
    Create a slide with lyrics
    """
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add text box
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(5.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = 1  # Middle alignment

    # Add the lyrics chunk
    for idx, line in enumerate(lyrics_chunk):
        if idx == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(28)
        p.font.name = 'Arial'
        p.font.color.rgb = RGBColor(0, 0, 0)

    # Set slide background to white
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

def create_presentation(all_lyrics: List[str], song_names: List[str], use_church_template: bool = True) -> str:
    """
    Create a PowerPoint presentation from multiple songs' lyrics.
    Lyrics are already intelligently grouped by Grok with ---SLIDE--- separators.

    Args:
        all_lyrics: List of lyrics for each song
        song_names: List of song titles
        use_church_template: If True, use church style (black bg, white text); if False, use default style
    """

    # Use church template if requested
    if use_church_template:
        try:
            from church_template import create_church_presentation_v3
            return create_church_presentation_v3(all_lyrics, song_names)
        except ImportError:
            print("Church template not available, using default template")
            # Fall through to default template
        except Exception as e:
            print(f"Error using church template: {e}, falling back to default")
            # Fall through to default template

    # Default template
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Process each song
    for idx, song_lyrics in enumerate(all_lyrics):
        song_name = song_names[idx] if idx < len(song_names) else f"Song {idx + 1}"

        # Add title slide for this song
        create_title_slide(prs, song_name)

        # Parse the pre-grouped lyrics from Grok
        lyric_groups = parse_grok_grouped_lyrics(song_lyrics)

        # Create a slide for each group
        for group in lyric_groups:
            if group:  # Only create slide if group has content
                create_lyrics_slide(prs, group)

    # Save the presentation
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"lyrics_presentation_{timestamp}.pptx"

    prs.save(filename)
    return filename
