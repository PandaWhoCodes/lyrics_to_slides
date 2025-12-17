"""
Church Template V3 - Working approach that preserves all formatting
Instead of copying slides, we modify the template directly
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
import os
import shutil
from copy import deepcopy


def insert_blank_separator(prs, separator_template):
    """
    Insert a blank black separator slide using template formatting.
    Copies the separator template and removes all text shapes.
    """
    template_layout = separator_template.slide_layout
    new_slide = prs.slides.add_slide(template_layout)

    # Remove all shapes to make it blank
    for shape in list(new_slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    return new_slide


def condense_repeated_lines(slide_text):
    """
    Condense consecutive repeated content within a slide.
    Handles both:
    1. Single lines repeated (e.g., "Hallelujah" x4)
    2. Multi-line blocks repeated (e.g., a 4-line verse repeated twice)

    Example 1 (single line):
    Input:
        Enthan Yesuvae
        Enthan Yesuvae
        Enthan Yesuvae
        Enthan Yesuvae
    Output:
        Enthan Yesuvae (x4)

    Example 2 (block):
    Input:
        Line A
        Line B
        Line A
        Line B
    Output:
        Line A
        Line B
        (x2)
    """
    if not slide_text:
        return slide_text

    lines = [line for line in slide_text.split('\n') if line.strip()]
    if len(lines) <= 2:
        return slide_text

    # First, try to detect repeated blocks
    condensed = try_condense_blocks(lines)
    if condensed:
        return condensed

    # Fall back to single-line repetition detection
    return condense_single_lines(lines)


def try_condense_blocks(lines):
    """
    Detect repeating blocks within the slide and condense them.
    Can handle partial repetitions (block repeats but other content follows).
    Returns condensed text if any repeating blocks found, None otherwise.
    """
    n = len(lines)
    result_parts = []
    i = 0

    while i < n:
        # Try to find a repeating block starting at position i
        found_block = False

        # Try different block sizes (prefer larger blocks, but include single lines)
        for block_size in range(min((n - i) // 2, 10), 0, -1):
            remaining = n - i
            if remaining < block_size * 2:
                continue

            # Get the potential block
            block = [line.strip().lower() for line in lines[i:i + block_size]]

            # Count how many times this block repeats consecutively
            repeat_count = 1
            pos = i + block_size

            while pos + block_size <= n:
                next_block = [line.strip().lower() for line in lines[pos:pos + block_size]]
                if next_block == block:
                    repeat_count += 1
                    pos += block_size
                else:
                    break

            # If block repeats at least twice, condense it
            if repeat_count >= 2:
                original_block = lines[i:i + block_size]

                # Check if all lines in the block are identical
                # If so, prefer single-line format even if block_size > 1
                block_lines_normalized = [line.strip().lower() for line in original_block]
                all_identical = len(set(block_lines_normalized)) == 1

                if block_size == 1 or all_identical:
                    # Use inline format: "Line (xN)"
                    total_repeats = block_size * repeat_count
                    result_parts.append(f"{original_block[0]} (x{total_repeats})")
                    print(f"    🔄 Condensed '{original_block[0][:30]}...' repeated {total_repeats} times → (x{total_repeats})")
                else:
                    # For multi-line blocks with different lines, put (xN) on its own line
                    result_parts.append('\n'.join(original_block) + f'\n(x{repeat_count})')
                    print(f"    🔄 Condensed {block_size}-line block repeated {repeat_count} times → (x{repeat_count})")
                i = pos
                found_block = True
                break

        if not found_block:
            # No repeating block found at this position, keep the line as-is
            result_parts.append(lines[i])
            i += 1

    # Check if any condensing happened
    if any('(x' in part for part in result_parts):
        return '\n'.join(result_parts)

    return None


def condense_single_lines(lines):
    """
    Condense consecutive identical single lines.
    """
    condensed_lines = []
    i = 0

    while i < len(lines):
        current_line = lines[i].strip()
        if not current_line:
            condensed_lines.append(lines[i])
            i += 1
            continue

        # Count consecutive repetitions
        count = 1
        j = i + 1
        while j < len(lines) and lines[j].strip().lower() == current_line.lower():
            count += 1
            j += 1

        # If repeated more than twice, condense it
        if count > 2:
            condensed_lines.append(f"{current_line} (x{count})")
            print(f"    🔄 Condensed '{current_line[:30]}...' repeated {count} times → (x{count})")
        else:
            # Keep original lines (1 or 2 repetitions)
            for k in range(count):
                condensed_lines.append(lines[i + k])

        i = j

    return '\n'.join(condensed_lines)


def deduplicate_slides(slide_texts):
    """
    Smart deduplication of slide content.
    - Removes consecutive duplicate slides (keeps first occurrence)
    - Returns list of unique slide texts in order

    Example:
    Input:  ["Verse 1", "Chorus", "Chorus", "Verse 2", "Chorus"]
    Output: ["Verse 1", "Chorus", "Verse 2", "Chorus"]

    The pastor doesn't need to flip through identical slides back-to-back.
    """
    if not slide_texts:
        return []

    deduplicated = [slide_texts[0]]

    for slide_text in slide_texts[1:]:
        # Normalize for comparison (strip whitespace, lowercase)
        current_normalized = slide_text.strip().lower()
        prev_normalized = deduplicated[-1].strip().lower()

        # Only add if different from previous slide
        if current_normalized != prev_normalized:
            deduplicated.append(slide_text)
        else:
            print(f"    ⏭️ Skipping duplicate slide (same as previous)")

    return deduplicated


def create_church_presentation_v3(lyrics_list, song_names):
    """
    Create a church-styled presentation by modifying template directly.
    This preserves all original formatting.

    Smart features:
    - Respects ---SLIDE--- markers set by user/AI for slide breaks
    - Deduplicates consecutive identical slides (chorus repeated back-to-back)
    - Keeps the flow natural for worship leaders
    """
    template_path = os.path.join(os.path.dirname(__file__), '..', 'reference_template.pptx')

    # Create a working copy of the template
    temp_working_path = os.path.join(os.path.dirname(__file__), '..', 'temp_working.pptx')
    shutil.copy(template_path, temp_working_path)

    print("Loading template...")
    prs = Presentation(temp_working_path)

    print(f"Original template has {len(prs.slides)} slides")
    print()

    # Keep only the first 3 slides (title, welcome, lyrics template)
    print("Keeping first 3 slides as base...")
    slides_to_keep = 3
    while len(prs.slides) > slides_to_keep:
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[-1]

    print(f"Now have {len(prs.slides)} slides")
    print()

    # Update date on first slide
    print("Updating date on title slide...")
    date_str = datetime.now().strftime("%d %b'%y")

    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            # Look for date pattern with apostrophe or month names
            if "'" in text or any(month in text for month in ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']):
                # Replace the entire text frame to avoid duplication
                shape.text_frame.text = date_str
                print(f"  Updated date to: {date_str}")
                break

    print()
    print("Creating lyrics slides...")

    # Get the lyrics template slide (slide 3)
    lyrics_template_slide = prs.slides[2]

    # Store separator template (Slide 2 - Welcome slide)
    separator_template_slide = prs.slides[1]

    # Process each song
    for song_idx, (lyrics, song_name) in enumerate(zip(lyrics_list, song_names)):
        # Insert separator slide BEFORE songs 2, 3, 4, etc. (not before first song)
        if song_idx > 0:
            insert_blank_separator(prs, separator_template_slide)
            print(f"  Inserted separator slide before song {song_idx + 1}")

        # Clean song name: remove " by Artist" for display
        display_title = song_name.split(' by ')[0].strip()

        print(f"  Processing '{display_title}'...")

        # Parse lyrics - respect AI's ---SLIDE--- groupings
        # Split by ---SLIDE--- marker to get pre-grouped sections
        sections = lyrics.split('---SLIDE---')

        # Build slide texts from sections (respecting user/AI slide breaks)
        all_slide_texts = []
        for section in sections:
            # Clean section: remove ---TITLE---, ---LYRICS--- markers and empty lines
            section_lines = []
            for line in section.split('\n'):
                line = line.strip()
                if line and not line.startswith('---'):
                    section_lines.append(line)

            if not section_lines:
                continue

            # Each section becomes one slide (respecting ---SLIDE--- markers)
            all_slide_texts.append('\n'.join(section_lines))

        # Smart deduplication - remove consecutive identical slides
        unique_slide_texts = deduplicate_slides(all_slide_texts)

        print(f"    📊 {len(all_slide_texts)} raw slides → {len(unique_slide_texts)} unique slides")

        total_slides = len(unique_slide_texts)

        for slide_idx, slide_text in enumerate(unique_slide_texts):
            # Condense repeated lines within this slide (e.g., "Hallelujah" x4 → "Hallelujah (x4)")
            slide_text = condense_repeated_lines(slide_text)

            # Duplicate the lyrics template slide using its own layout
            # Use the same layout as the template to preserve structure
            template_layout = lyrics_template_slide.slide_layout
            new_slide = prs.slides.add_slide(template_layout)

            # Remove all existing shapes from the new slide
            for shape in list(new_slide.shapes):
                sp = shape.element
                sp.getparent().remove(sp)

            # Deep copy all shapes from template
            for shape in lyrics_template_slide.shapes:
                el = shape.element
                newel = deepcopy(el)
                new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

            # Now replace text in the new slide
            for shape in new_slide.shapes:
                if shape.has_text_frame:
                    # Get position to determine what type of text this is
                    # Title text (small text at top)
                    if shape.top < Inches(1):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip():
                                    run.text = display_title
                                    break
                    # Main lyrics text (large text in middle)
                    elif Inches(1) < shape.top < Inches(4.5):
                        # Clear existing text and set new lyrics
                        shape.text_frame.text = slide_text

                        # Remove bullet formatting from all paragraphs
                        for paragraph in shape.text_frame.paragraphs:
                            # Remove bullets by clearing the paragraph properties
                            if hasattr(paragraph, '_element') and paragraph._element is not None:
                                pPr = paragraph._element.get_or_add_pPr()
                                # Remove bullet character element if it exists
                                buChar = pPr.find('.//a:buChar', {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                                if buChar is not None:
                                    pPr.remove(buChar)
                                # Set buNone to explicitly have no bullets
                                from pptx.oxml.ns import qn
                                buNone = pPr.find(qn('a:buNone'))
                                if buNone is None:
                                    from pptx.oxml import parse_xml
                                    buNone = parse_xml('<a:buNone xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
                                    pPr.insert(0, buNone)
                    # Slide number (small text at bottom right)
                    elif shape.top > Inches(4.5) and shape.left > Inches(7):
                        current_slide_num = slide_idx + 1
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.text = f"{current_slide_num}/{total_slides}"
                                break

    # Delete slide 3 (the template slide) after we're done using it
    print()
    print("Removing template slide (slide 3)...")
    # Get the slide ID for slide 3 (index 2)
    slide_id = prs.slides._sldIdLst[2]
    prs.part.drop_rel(slide_id.rId)
    del prs.slides._sldIdLst[2]

    # Save presentation
    output_path = os.path.join(os.path.dirname(__file__), '..', 'church_style_output.pptx')
    prs.save(output_path)

    # Clean up temp file
    try:
        os.remove(temp_working_path)
    except:
        pass

    print()
    print(f"✅ Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    # Test with repeated chorus AND repeated lines within a slide
    test_lyrics = [
        """---TITLE---
Amazing Grace
---LYRICS---
Amazing grace, how sweet the sound
That saved a wretch like me
I once was lost but now am found
Was blind, but now I see
---SLIDE---
Hallelujah
Hallelujah
Hallelujah
Hallelujah
---SLIDE---
'Twas grace that taught my heart to fear
And grace my fears relieved
How precious did that grace appear
The hour I first believed
---SLIDE---
Enthan Yesuvae
Enthan Yesuvae
Enthan Yesuvae
Enthan Yesuvae
---SLIDE---
Through many dangers, toils and snares
I have already come
'Tis grace hath brought me safe thus far
And grace will lead me home"""
    ]
    test_songs = ["Amazing Grace"]

    print("Testing smart deduplication and line condensing...")
    print("Features being tested:")
    print("  1. Consecutive duplicate slides are removed")
    print("  2. Repeated lines within a slide (>2x) become 'Line (xN)'")
    print()
    print("Expected:")
    print("  - 'Hallelujah' x4 → 'Hallelujah (x4)'")
    print("  - 'Enthan Yesuvae' x4 → 'Enthan Yesuvae (x4)'")
    print()

    create_church_presentation_v3(test_lyrics, test_songs)
